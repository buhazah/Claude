"""Format extractors.

Each turns a source into located blocks. The locator is the whole point: a PDF
page, a slide number, a spreadsheet cell range, a line span. Without it a
citation degrades to "it's in this file somewhere", which is not an answer.

Third-party parsers are imported *inside* their extractor, so a missing or
broken dependency degrades one format rather than preventing the module from
loading — and ingestion reports which format it could not read.
"""

from __future__ import annotations

import csv
import io
import re
from pathlib import Path
from typing import Any

import structlog

from jarvis.knowledge.models import Block, Extracted, SourceKind, kind_for

log = structlog.get_logger(__name__)

MAX_SHEET_ROWS = 2000
MAX_CODE_LINES_PER_BLOCK = 80


class ExtractionError(Exception):
    """A source could not be read. Carries why, for the user."""


def extract_path(path: Path) -> Extracted:
    """Read a file into located blocks, choosing an extractor by suffix."""
    kind = kind_for(path.suffix)
    extractor = _EXTRACTORS.get(kind)
    if extractor is None:
        # Unknown suffixes are attempted as text: many useful files (LICENSE,
        # Dockerfile, .env.example) have no recognised extension but are plain.
        return _extract_text(path, SourceKind.TEXT)
    return extractor(path, kind)


def _read_text(path: Path) -> str:
    return path.read_bytes().decode("utf-8", errors="replace")


def _extract_text(path: Path, kind: SourceKind) -> Extracted:
    text = _read_text(path)
    blocks = [
        Block(text=part.strip(), locator=f"¶{i + 1}", order=i)
        for i, part in enumerate(re.split(r"\n\s*\n", text))
        if part.strip()
    ]
    return Extracted(title=path.name, kind=kind, blocks=blocks)


def _extract_markdown(path: Path, kind: SourceKind) -> Extracted:
    """Split on headings, so a locator names the section a reader would look for."""
    text = _read_text(path)
    blocks: list[Block] = []
    heading = ""
    buffer: list[str] = []

    def flush() -> None:
        body = "\n".join(buffer).strip()
        if body:
            blocks.append(Block(text=body, locator=heading or "intro", order=len(blocks)))
        buffer.clear()

    for line in text.splitlines():
        if match := re.match(r"^(#{1,6})\s+(.*)", line):
            flush()
            heading = match.group(2).strip()
            buffer.append(line)
        else:
            buffer.append(line)
    flush()

    title = next(
        (
            m.group(1).strip()
            for m in [re.match(r"^#\s+(.*)", ln) for ln in text.splitlines()[:5]]
            if m
        ),
        path.name,
    )
    return Extracted(title=title, kind=kind, blocks=blocks)


def _extract_pdf(path: Path, kind: SourceKind) -> Extracted:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - dependency is declared
        raise ExtractionError("pypdf is not installed") from exc

    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        raise ExtractionError(f"could not read PDF: {exc}") from exc

    blocks = [
        Block(text=text.strip(), locator=f"p. {number}", order=number - 1)
        for number, page in enumerate(reader.pages, start=1)
        if (text := (page.extract_text() or "")) and text.strip()
    ]
    info: dict[str, Any] = dict(reader.metadata or {})
    title = str(info.get("/Title") or "").strip() or path.name
    return Extracted(
        title=title,
        kind=kind,
        blocks=blocks,
        metadata={"pages": str(len(reader.pages))},
    )


def _extract_docx(path: Path, kind: SourceKind) -> Extracted:
    try:
        import docx
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("python-docx is not installed") from exc

    try:
        document = docx.Document(str(path))
    except Exception as exc:
        raise ExtractionError(f"could not read DOCX: {exc}") from exc

    blocks: list[Block] = []
    heading = ""
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name if paragraph.style else ""
        if style.startswith("Heading"):
            heading = text
        blocks.append(Block(text=text, locator=heading or "body", order=len(blocks)))

    # Tables carry the numbers people actually search for; flatten them rather
    # than dropping them.
    for index, table in enumerate(document.tables, start=1):
        rows = [
            " | ".join(cell.text.strip() for cell in row.cells)
            for row in table.rows
            if any(cell.text.strip() for cell in row.cells)
        ]
        if rows:
            blocks.append(Block(text="\n".join(rows), locator=f"table {index}", order=len(blocks)))

    title = next(
        (
            p.text.strip()
            for p in document.paragraphs
            if p.style and p.style.name == "Title" and p.text.strip()
        ),
        path.name,
    )
    return Extracted(title=title, kind=kind, blocks=blocks)


def _extract_pptx(path: Path, kind: SourceKind) -> Extracted:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("python-pptx is not installed") from exc

    try:
        deck = Presentation(str(path))
    except Exception as exc:
        raise ExtractionError(f"could not read PPTX: {exc}") from exc

    blocks: list[Block] = []
    for number, slide in enumerate(deck.slides, start=1):
        parts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        notes = getattr(slide, "notes_slide", None)
        # Speaker notes carry the argument the slide only gestures at.
        if notes and (note_text := notes.notes_text_frame.text.strip()):
            parts.append(f"[speaker notes] {note_text}")
        if parts:
            blocks.append(Block(text="\n".join(parts), locator=f"slide {number}", order=number - 1))

    return Extracted(
        title=path.stem, kind=kind, blocks=blocks, metadata={"slides": str(len(deck.slides))}
    )


def _extract_xlsx(path: Path, kind: SourceKind) -> Extracted:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:  # pragma: no cover
        raise ExtractionError("openpyxl is not installed") from exc

    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        raise ExtractionError(f"could not read XLSX: {exc}") from exc

    blocks: list[Block] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        first_row = 0
        for index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            if index > MAX_SHEET_ROWS:
                rows.append(f"[truncated after {MAX_SHEET_ROWS} rows]")
                break
            cells = [str(v) for v in row if v is not None]
            if not cells:
                continue
            first_row = first_row or index
            rows.append(" | ".join(cells))
        if rows:
            span = f"{sheet.title}!{first_row}:{first_row + len(rows) - 1}"
            blocks.append(Block(text="\n".join(rows), locator=span, order=len(blocks)))
    workbook.close()

    return Extracted(title=path.stem, kind=kind, blocks=blocks)


def _extract_csv(path: Path, kind: SourceKind) -> Extracted:
    text = _read_text(path)
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows = [" | ".join(row) for row in reader if any(cell.strip() for cell in row)]

    header = rows[0] if rows else ""
    blocks: list[Block] = []
    # Repeat the header in each block: a row of numbers without its column
    # names is unusable once retrieved out of context.
    for start in range(1, len(rows), 50):
        window = rows[start : start + 50]
        blocks.append(
            Block(
                text="\n".join([header, *window]),
                locator=f"rows {start + 1}–{start + len(window)}",
                order=len(blocks),
            )
        )
    return Extracted(title=path.name, kind=kind, blocks=blocks)


def _extract_html(path: Path, kind: SourceKind) -> Extracted:
    return Extracted(title=path.name, kind=kind, blocks=html_blocks(_read_text(path)))


def _extract_code(path: Path, kind: SourceKind) -> Extracted:
    """Split code on line spans, so a citation points at lines you can open."""
    lines = _read_text(path).splitlines()
    blocks: list[Block] = []
    for start in range(0, len(lines), MAX_CODE_LINES_PER_BLOCK):
        window = lines[start : start + MAX_CODE_LINES_PER_BLOCK]
        if not any(line.strip() for line in window):
            continue
        blocks.append(
            Block(
                text="\n".join(window),
                locator=f"lines {start + 1}–{start + len(window)}",
                order=len(blocks),
            )
        )
    return Extracted(title=path.name, kind=kind, blocks=blocks, metadata={"lines": str(len(lines))})


_SCRIPT_RE = re.compile(r"<(script|style)\b.*?</\1>", re.S | re.I)
_TAG_RE = re.compile(r"<[^>]+>")
_TITLE_RE = re.compile(r"<title[^>]*>(.*?)</title>", re.S | re.I)
_ENTITIES = {"&nbsp;": " ", "&amp;": "&", "&lt;": "<", "&gt;": ">", "&quot;": '"', "&#39;": "'"}


def html_title(html: str) -> str:
    match = _TITLE_RE.search(html)
    return match.group(1).strip() if match else ""


def html_blocks(html: str) -> list[Block]:
    """Strip markup to readable paragraphs.

    Deliberately not a full parser: the goal is retrievable prose, and a
    dependency-free strip is predictable and fast. Headings become locators.
    """
    body = _SCRIPT_RE.sub(" ", html)
    # Keep heading text as a locator before the tags are stripped away.
    body = re.sub(r"<h[1-6][^>]*>(.*?)</h[1-6]>", r"\n\n§\1§\n\n", body, flags=re.S | re.I)
    body = re.sub(r"<(p|div|br|li|tr)\b[^>]*>", "\n", body, flags=re.I)
    body = _TAG_RE.sub(" ", body)
    for entity, char in _ENTITIES.items():
        body = body.replace(entity, char)

    blocks: list[Block] = []
    heading = ""
    for part in re.split(r"\n\s*\n", body):
        text = re.sub(r"[ \t]+", " ", part).strip()
        if not text:
            continue
        if text.startswith("§") and text.endswith("§"):
            heading = text.strip("§").strip()
            continue
        blocks.append(Block(text=text, locator=heading or "body", order=len(blocks)))
    return blocks


def extract_html_document(html: str, *, source: str, title: str = "") -> Extracted:
    return Extracted(
        title=title or html_title(html) or source,
        kind=SourceKind.URL,
        blocks=html_blocks(html),
        metadata={"source": source},
    )


_EXTRACTORS = {
    SourceKind.TEXT: _extract_text,
    SourceKind.MARKDOWN: _extract_markdown,
    SourceKind.PDF: _extract_pdf,
    SourceKind.DOCX: _extract_docx,
    SourceKind.PPTX: _extract_pptx,
    SourceKind.XLSX: _extract_xlsx,
    SourceKind.CSV: _extract_csv,
    SourceKind.HTML: _extract_html,
    SourceKind.CODE: _extract_code,
}
